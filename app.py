
import streamlit as st
from pptx import Presentation
import fitz
import tempfile
import re
from datetime import datetime

from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, PageBreak
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import A4

# =========================================================
# 页面
# =========================================================
st.title("📊 Investment Product Explorer V3")
st.caption("Clean Classification + Exportable Report")

report_date = st.date_input("📅 报告日期", value=datetime.today())

st.info("""
上传同一份文件：
1️⃣ PPTX（分类）
2️⃣ PDF（展示）

⚠️ 页数必须一致
""")

ppt_file = st.file_uploader("Upload PPTX", type=["pptx"])
pdf_file = st.file_uploader("Upload PDF", type=["pdf"])

slides_data = []
image_list = []

# =========================================================
# PDF → 图片
# =========================================================
def pdf_to_images(pdf_file):

    images = []

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(pdf_file.read())
        tmp_path = tmp.name

    pdf = fitz.open(tmp_path)

    for i in range(len(pdf)):
        page = pdf[i]
        pix = page.get_pixmap()
        img_path = f"/tmp/page_{i+1}.png"
        pix.save(img_path)
        images.append(img_path)

    return images

# =========================================================
# 文本处理
# =========================================================
def clean_text(text):
    return re.sub(r"\s+", " ", text.lower())

# =========================================================
# ✅ V3 分类逻辑（最终稳定版）
# =========================================================
def classify(text):

    # ✅ DCI优先（核心规则）
    if "dci" in text:

        if "dual" in text and "range accrual" in text:
            return "DCI", "Dual Range DCI"

        elif "range accrual" in text:
            return "DCI", "Range Accrual DCI"

        else:
            return "DCI", "Vanilla DCI"

    # ✅ Accrual Note（纯）
    elif "range accrual" in text:

        if "dual" in text:
            return "Accrual Note", "Dual Accrual Note"
        else:
            return "Accrual Note", "Accrual Note"

    # ✅ 主类
    elif "fcn" in text:
        return "FCN", "FCN"

    elif "sharkfin" in text:
        return "Sharkfin", "Sharkfin"

    elif "aq" in text:
        return "AQ", "Accumulator"

    elif "fund" in text or "bluebay" in text or "singularity" in text:
        return "Fund", "Fund"

    # ✅ Others
    elif "twinwin" in text:
        return "Others", "Twinwin"

    elif "dual digital" in text or "warrant" in text:
        return "Others", "Dual Digital"

    elif "ben" in text:
        return "Others", "BEN"

    elif "dq" in text:
        return "Others", "DQ"

    elif "tarf" in text:
        return "Others", "TARF"

    elif "inverse floater" in text:
        return "Others", "Inverse Floater"

    elif "stable note" in text:
        return "Others", "Stable Note"

    else:
        return "Others", "Unknown"

# =========================================================
# ✅ PDF（带目录）
# =========================================================
def generate_toc_pdf(grouped, image_list):

    output_path = "/tmp/report_v3.pdf"

    doc = SimpleDocTemplate(output_path, pagesize=A4)
    styles = getSampleStyleSheet()
    content = []

    ordered_main = [
        "FCN",
        "Accrual Note",
        "DCI",
        "Sharkfin",
        "AQ",
        "Fund",
        "Others"
    ]

    # ✅ 封面
    content.append(Paragraph("Investment Product Report", styles["Title"]))
    content.append(Spacer(1, 50))
    content.append(Paragraph(str(report_date), styles["Normal"]))
    content.append(PageBreak())

    # ✅ 目录
    content.append(Paragraph("Table of Contents", styles["Heading1"]))
    content.append(Spacer(1, 20))

    page_index = 2
    page_map = {}

    for main in ordered_main:
        if main in grouped:
            page_map[main] = page_index
            page_index += sum(len(v) for v in grouped[main].values()) + 1

    for main in ordered_main:
        if main in page_map:
            content.append(
                Paragraph(f"{main} ...... Page {page_map[main]}", styles["Normal"])
            )
            content.append(Spacer(1, 8))

    content.append(PageBreak())

    # ✅ 内容
    for main in ordered_main:

        if main not in grouped:
            continue

        content.append(Paragraph(f"<b>{main}</b>", styles["Heading1"]))
        content.append(Spacer(1, 10))

        # Others（有子类）
        if main == "Others":

            sorted_sub = sorted(
                grouped[main].keys(),
                key=lambda x: (x == "Unknown", x)
            )

            for sub in sorted_sub:

                content.append(Paragraph(sub, styles["Heading2"]))

                for slide in grouped[main][sub]:
                    page_num = slide["page"]
                    img = image_list[page_num - 1]

                    content.append(Paragraph(f"Page {page_num}", styles["Normal"]))
                    content.append(RLImage(img, width=400, height=300))
                    content.append(Spacer(1, 12))

        else:

            for sub in grouped[main].values():
                for slide in sub:
                    page_num = slide["page"]
                    img = image_list[page_num - 1]

                    content.append(Paragraph(f"Page {page_num}", styles["Normal"]))
                    content.append(RLImage(img, width=400, height=300))
                    content.append(Spacer(1, 12))

        content.append(PageBreak())

    doc.build(content)

    return output_path

# =========================================================
# 主逻辑
# =========================================================
if ppt_file and pdf_file:

    prs = Presentation(ppt_file)

    for i, slide in enumerate(prs.slides):

        text_content = ""

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content += shape.text + " "

        slides_data.append({
            "page": i + 1,
            "text": text_content
        })

    image_list = pdf_to_images(pdf_file)

    if len(slides_data) != len(image_list):
        st.error("❌ 页数不一致")
        st.stop()

    grouped = {}

    for slide in slides_data:

        text = clean_text(slide["text"])
        main, sub = classify(text)

        if main not in grouped:
            grouped[main] = {}

        if sub not in grouped[grouped[main]]:
            grouped[main][sub] = []

        grouped[main][sub].append(slide)

    # ✅ UI
    ordered_main = [
        "FCN",
        "Accrual Note",
        "DCI",
        "Sharkfin",
        "AQ",
        "Fund",
        "Others"
    ]

    for main in ordered_main:

        if main not in grouped:
            continue

        count = sum(len(v) for v in grouped[main].values())

        with st.expander(f"📂 {main} ({count})"):

            if main == "Others":

                sorted_sub = sorted(
                    grouped[main].keys(),
                    key=lambda x: (x == "Unknown", x)
                )

                for sub in sorted_sub:

                    slides = grouped[main][sub]

                    with st.expander(f"{sub} ({len(slides)})"):

                        for s in slides:
                            st.markdown(f"**Page {s['page']}**")
                            st.image(image_list[s["page"] - 1])

            else:

                all_slides = []
                for v in grouped[main].values():
                    all_slides.extend(v)

                for s in all_slides:
                    st.markdown(f"**Page {s['page']}**")
                    st.image(image_list[s["page"] - 1])

    # ✅ PDF导出
    pdf_path = generate_toc_pdf(grouped, image_list)

    with open(pdf_path, "rb") as f:
        st.download_button("📄 下载PDF报告（带目录）", f, "report.pdf")

